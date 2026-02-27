#!/usr/bin/env python3
"""
Aurelian Manufacturing — Customer-Segment PPTX Generator
Generates a completely new customer-facing presentation targeting
Defense, Energy, and Maritime/Infrastructure segments.

Design Manual v2.0 compliant. All measurements in inches.
"""

import base64
import io
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── PATHS ──────────────────────────────────────────────────────────────────────
DESIGN_MANUAL_DIR = Path(r"C:\Users\mrntau\OneDrive - Mosseregionens Næringsutvikling AS\Skrivebord\Design MAnual")
OUTPUT_DIR = Path(r"C:\Users\mrntau\OneDrive - Mosseregionens Næringsutvikling AS\Skrivebord\Aurelian\Aurelian Manufacturing\Aurelian_VDR\07_Presentations")
OUTPUT_FILE = OUTPUT_DIR / "Aurelian_Customer_Presentation_v7.pptx"

# ── DESIGN TOKENS ──────────────────────────────────────────────────────────────
RED = RGBColor(0xF5, 0x05, 0x37)       # Aurelian Red / Pantone 032 C
BLACK = RGBColor(0x2B, 0x2B, 0x2B)     # Near-black text
PURE_BLACK = RGBColor(0x00, 0x00, 0x00) # Cover/closing backgrounds
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)      # Secondary text
PANEL_GRAY = RGBColor(0xF5, 0xF5, 0xF5) # Panel backgrounds
DIVIDER_GRAY = RGBColor(0xD3, 0xD3, 0xD3) # Borders

FONT_BODY = "Calibri"
FONT_CHART = "Arial"

SLIDE_WIDTH = Inches(10.0)
SLIDE_HEIGHT = Inches(5.625)
MARGIN = Inches(0.5)
CONTENT_WIDTH = Inches(9.0)

# ── HELPERS ────────────────────────────────────────────────────────────────────

def load_b64_image(filename):
    """Load a base64-encoded image from Design Manual dir, return BytesIO."""
    filepath = DESIGN_MANUAL_DIR / filename
    with open(filepath, "r") as f:
        data = f.read().strip()
    return io.BytesIO(base64.b64decode(data))


def set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                font_color=BLACK, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    """Add a text box with a single run of formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Set vertical anchor
    try:
        txBox.text_frame._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"))
    except:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox


def add_multi_text(slide, left, top, width, height, parts, alignment=PP_ALIGN.LEFT):
    """Add text box with multiple formatted runs.
    parts: list of dicts with keys: text, size, color, bold, italic, font_name
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    for part in parts:
        run = p.add_run()
        run.text = part.get("text", "")
        run.font.size = Pt(part.get("size", 14))
        run.font.color.rgb = part.get("color", BLACK)
        run.font.bold = part.get("bold", False)
        run.font.italic = part.get("italic", False)
        run.font.name = part.get("font_name", FONT_BODY)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    font_color=BLACK, spacing_pt=6):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(spacing_pt)
        p.level = 0
        # Bullet character
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = FONT_BODY
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # No border by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        if border_width:
            shape.line.width = Pt(border_width)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_circle(slide, left, top, diameter, fill_color=RED):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_footer(slide, bg_is_dark=False):
    """Add CONFIDENTIAL + date footer to slide."""
    text_color = GRAY
    # Confidential left
    add_textbox(slide, MARGIN, Inches(5.2), Inches(2.0), Inches(0.3),
                "CONFIDENTIAL", font_size=10, font_color=text_color)
    # Date right
    add_textbox(slide, Inches(7.5), Inches(5.2), Inches(2.0), Inches(0.3),
                "16 February 2026", font_size=10, font_color=text_color,
                alignment=PP_ALIGN.RIGHT)


def add_slide_title(slide, title, subtitle=None):
    """Add standard slide title (and optional subtitle)."""
    add_textbox(slide, MARGIN, Inches(0.3), CONTENT_WIDTH, Inches(0.6),
                title, font_size=28, font_color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, MARGIN, Inches(0.85), CONTENT_WIDTH, Inches(0.4),
                    subtitle, font_size=16, font_color=RED, bold=True, italic=True)


def create_card(slide, left, top, width, height, title, body, title_size=14, body_size=11):
    """Create a panel/card with gray background."""
    add_rect(slide, left, top, width, height, fill_color=PANEL_GRAY)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35),
                title, font_size=title_size, font_color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55),
                body, font_size=body_size, font_color=BLACK, alignment=PP_ALIGN.CENTER)


def create_stat_badge(slide, left, top, width, height, number, label):
    """Create a stat badge: white number on red background."""
    add_rounded_rect(slide, left, top, width, height, fill_color=RED)
    add_textbox(slide, left, top + Inches(0.05), width, Inches(0.5),
                number, font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + height + Inches(0.05), width, Inches(0.3),
                label, font_size=10, font_color=GRAY, alignment=PP_ALIGN.CENTER)


def create_table(slide, left, top, width, rows_data, col_widths=None):
    """Create a styled table.
    rows_data: list of lists. First row is header.
    col_widths: list of Inches values.
    """
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    row_height = Inches(0.35)
    table_height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = FONT_BODY
            run.font.size = Pt(10)

            if row_idx == 0:
                # Header row
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.size = Pt(10)
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLACK
            else:
                run.font.color.rgb = BLACK
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PANEL_GRAY
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

            # Remove all borders, then add horizontal only
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                for existing in tcPr.findall(qn(border_name)):
                    tcPr.remove(existing)

    return table_shape


def add_numbered_step(slide, left, top, number, title, description):
    """Add a numbered step with red circle marker."""
    circle = add_circle(slide, left, top, Inches(0.35), fill_color=RED)
    add_textbox(slide, left + Inches(0.02), top + Inches(0.02), Inches(0.35), Inches(0.35),
                str(number), font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.5), top - Inches(0.02), Inches(3.5), Inches(0.3),
                title, font_size=13, font_color=BLACK, bold=True)
    add_textbox(slide, left + Inches(0.5), top + Inches(0.25), Inches(3.5), Inches(0.35),
                description, font_size=10, font_color=GRAY)


def add_callout_box(slide, left, top, width, height, text, font_size=11):
    """Add a callout box with red left border."""
    # Red left border line
    add_rect(slide, left, top, Inches(0.04), height, fill_color=RED)
    # Gray background panel
    add_rect(slide, left + Inches(0.04), top, width - Inches(0.04), height, fill_color=PANEL_GRAY)
    # Text
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.35), height - Inches(0.2),
                text, font_size=font_size, font_color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def build_slide_01_cover(prs, logo_img):
    """COVER SLIDE — Black background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, PURE_BLACK)

    # Logo
    logo_img.seek(0)
    slide.shapes.add_picture(logo_img, Inches(4.25), Inches(0.5), Inches(1.5), Inches(1.5))

    # Company name
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(9.0), Inches(0.8),
                "AURELIAN MANUFACTURING", font_size=42, font_color=WHITE,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Red decorative bar
    add_rect(slide, Inches(3.5), Inches(3.05), Inches(3.0), Inches(0.04), fill_color=RED)

    # Tagline
    add_textbox(slide, Inches(1.0), Inches(3.2), Inches(8.0), Inches(0.5),
                "Production as a Service", font_size=22,
                font_color=RED, italic=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(1.0), Inches(3.7), Inches(8.0), Inches(0.4),
                "Customer Capabilities Presentation", font_size=16,
                font_color=WHITE, alignment=PP_ALIGN.CENTER)

    add_footer(slide, bg_is_dark=True)


def build_slide_02_problem(prs):
    """THE PROBLEM — White + dark panel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "The Problem: Uncontrollable Supply Chains")

    # Left content — bullet points
    items = [
        "39,000 skilled worker shortage in Norway (NAV 2025)",
        "Defense, energy, and maritime sectors facing critical delivery delays",
        "Traditional job shops operating at 20-45% utilization",
        "Average 2.5 FTE per CNC machine \u2014 labor-intensive, inflexible",
    ]
    add_bullet_list(slide, MARGIN, Inches(1.35), Inches(5.5), Inches(2.5), items, font_size=13)

    # Right dark panel
    panel_left = Inches(6.3)
    panel_width = Inches(3.4)
    add_rect(slide, panel_left, Inches(1.2), panel_width, Inches(3.5), fill_color=BLACK)

    add_textbox(slide, panel_left + Inches(0.2), Inches(1.35), Inches(3.0), Inches(0.35),
                "THE RESULT", font_size=16, font_color=WHITE, bold=True)

    # Stat 1
    add_textbox(slide, panel_left + Inches(0.2), Inches(1.85), Inches(3.0), Inches(0.35),
                "12-18 months", font_size=24, font_color=WHITE, bold=True)
    add_textbox(slide, panel_left + Inches(0.2), Inches(2.2), Inches(3.0), Inches(0.25),
                "Defense procurement lead times", font_size=10, font_color=RED)

    # Stat 2
    add_textbox(slide, panel_left + Inches(0.2), Inches(2.65), Inches(3.0), Inches(0.35),
                "Supply risk", font_size=20, font_color=RED, bold=True)
    add_textbox(slide, panel_left + Inches(0.2), Inches(3.0), Inches(3.0), Inches(0.25),
                "threatens critical infrastructure delivery", font_size=10, font_color=WHITE)

    # Stat 3
    add_textbox(slide, panel_left + Inches(0.2), Inches(3.45), Inches(3.0), Inches(0.35),
                "HMLV bottleneck", font_size=20, font_color=WHITE, bold=True)
    add_textbox(slide, panel_left + Inches(0.2), Inches(3.8), Inches(3.0), Inches(0.35),
                "High-Mix, Low-Volume capacity crisis across all sectors", font_size=10, font_color=GRAY)

    add_footer(slide)


def build_slide_03_solution(prs, icon_gears, icon_robot, icon_check):
    """OUR SOLUTION — Three value proposition cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Controllable Capacity in Uncontrollable Supply Chains",
                    "Autonomous 5-axis CNC manufacturing designed for delivery reliability")

    # Three cards
    card_w = Inches(2.7)
    card_h = Inches(2.8)
    card_y = Inches(1.5)
    gap = Inches(0.45)
    start_x = MARGIN

    cards = [
        ("Supply Security", "Guaranteed capacity, Norwegian-based, European trusted framework. Regional, transparent supply chains.", icon_gears),
        ("Autonomous Operations", "24/7 autonomous scheduling, 0.8 FTE/CNC vs industry 2.5. Sub-linear staffing model.", icon_robot),
        ("Absolute Reliability", "Verifiable quality traceability, digital thread integration, predictable delivery at every stage.", icon_check),
    ]

    for i, (title, body, icon) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h, fill_color=PANEL_GRAY)

        # Icon
        if icon:
            icon.seek(0)
            slide.shapes.add_picture(icon, x + Inches(1.0), card_y + Inches(0.2), Inches(0.7), Inches(0.7))

        add_textbox(slide, x + Inches(0.15), card_y + Inches(1.0), card_w - Inches(0.3), Inches(0.35),
                    title, font_size=14, font_color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), card_y + Inches(1.4), card_w - Inches(0.3), Inches(1.3),
                    body, font_size=11, font_color=BLACK, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def build_slide_04_segments(prs):
    """CUSTOMER SEGMENTS — Three tier cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Customer Segments",
                    "Designed for high-precision, high-reliability sectors")

    card_w = Inches(2.7)
    card_h = Inches(2.8)
    card_y = Inches(1.5)
    gap = Inches(0.45)

    segments = [
        ("TIER 1: DEFENSE & AEROSPACE",
         "Kongsberg \u2022 Nammo \u2022 Saab \u2022 BAE Systems",
         "AQAP 2110 / AS9100 required\n12-18 month sales cycle",
         RED),
        ("TIER 2: ENERGY & OFFSHORE",
         "Equinor \u2022 Aker Solutions \u2022 Vattenfall",
         "NORSOK-compatible\n6-12 month sales cycle",
         BLACK),
        ("TIER 3: MARITIME & INFRA",
         "Vard \u2022 Ulstein \u2022 Bane NOR",
         "ISO 9001 sufficient\n3-6 month sales cycle",
         GRAY),
    ]

    for i, (title, customers, reqs, accent) in enumerate(segments):
        x = MARGIN + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h, fill_color=PANEL_GRAY)

        # Accent bar at top of card
        add_rect(slide, x, card_y, card_w, Inches(0.06), fill_color=accent)

        add_textbox(slide, x + Inches(0.15), card_y + Inches(0.2), card_w - Inches(0.3), Inches(0.5),
                    title, font_size=12, font_color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), card_y + Inches(0.8), card_w - Inches(0.3), Inches(0.8),
                    customers, font_size=11, font_color=BLACK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), card_y + Inches(1.8), card_w - Inches(0.3), Inches(0.8),
                    reqs, font_size=10, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def build_slide_05_defense(prs):
    """DEFENSE DRIVERS."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Defense: Unprecedented Demand Surge",
                    "NATO commitments, technological sovereignty, HMLV bottleneck")

    # Numbered drivers (left column)
    drivers = [
        ("NATO 5% GDP target by 2035", "Structural, multi-decade increase in defense spending across all alliance members"),
        ("Norway\u2019s 1,624B NOK defense plan", "12-year commitment (2025-2036) nearly doubling annual defense budgets"),
        ("EU 50% EU-sourced defense by 2030", "EDIS targets create reshoring imperative for European manufacturing"),
        ("HMLV manufacturing bottleneck", "Existing suppliers at capacity \u2014 critical gap for new certified production"),
    ]

    for i, (title, desc) in enumerate(drivers):
        y = Inches(1.45) + i * Inches(0.75)
        add_numbered_step(slide, MARGIN, y, i + 1, title, desc)

    # Right stat badge
    create_stat_badge(slide, Inches(6.8), Inches(1.5), Inches(2.5), Inches(0.65),
                      "1,624B NOK", "Norway defense plan 2025-2036")

    # Callout box
    add_callout_box(slide, Inches(6.5), Inches(3.0), Inches(3.0), Inches(1.2),
                    "Critical need: Trusted, certified, Norwegian-based manufacturing capacity for defense components",
                    font_size=10)

    add_footer(slide)


def build_slide_06_energy(prs):
    """ENERGY DRIVERS."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Energy: Europe\u2019s Critical Infrastructure",
                    "Dual-track O&G + energy transition, NORSOK, supply security")

    drivers = [
        ("Norway as Europe\u2019s most critical energy supplier", "Strategic role amplified post-Ukraine \u2014 energy security is national security"),
        ("Dual-track: O&G maintenance + energy transition", "Offshore wind, CCS, hydrogen \u2014 all require precision-machined components"),
        ("NORSOK compliance required", "Norwegian/Nordic standards for safety-critical energy components"),
        ("Supply chain resilience imperative", "Nearshoring and regional sourcing to reduce geopolitical supply risk"),
    ]

    for i, (title, desc) in enumerate(drivers):
        y = Inches(1.45) + i * Inches(0.75)
        add_numbered_step(slide, MARGIN, y, i + 1, title, desc)

    # Callout box
    add_callout_box(slide, Inches(6.5), Inches(2.5), Inches(3.0), Inches(1.5),
                    "Energy customers require: NORSOK-compatible processes, verifiable traceability, absolute delivery reliability for safety-critical components",
                    font_size=10)

    add_footer(slide)


def build_slide_07_maritime(prs):
    """MARITIME & INFRASTRUCTURE DRIVERS."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Maritime & Infrastructure: Nearshoring & Precision",
                    "Regional secure supply chains, next-gen autonomous vessels")

    drivers = [
        ("Reshoring / nearshoring trend", "Regional supply chain security after global disruptions \u2014 proximity matters"),
        ("High-precision autonomous vessels", "Next-gen maritime systems demand tighter tolerances and digital traceability"),
        ("Infrastructure modernization", "Rail, energy, transport \u2014 predictable delivery windows are critical"),
    ]

    for i, (title, desc) in enumerate(drivers):
        y = Inches(1.45) + i * Inches(0.8)
        add_numbered_step(slide, MARGIN, y, i + 1, title, desc)

    # Dark stat panel (right)
    add_rect(slide, Inches(6.5), Inches(1.5), Inches(3.0), Inches(1.2), fill_color=BLACK)
    add_textbox(slide, Inches(6.5), Inches(1.65), Inches(3.0), Inches(0.5),
                "3-6 months", font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.5), Inches(2.15), Inches(3.0), Inches(0.3),
                "typical sales cycle", font_size=11, font_color=RED, alignment=PP_ALIGN.CENTER)

    add_callout_box(slide, Inches(6.5), Inches(3.2), Inches(3.0), Inches(1.0),
                    "Maritime/infrastructure customers value: Fast onboarding, predictable capacity, regional proximity",
                    font_size=10)

    add_footer(slide)


def build_slide_08_facility(prs):
    """FACILITY — Purpose-built for autonomy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Purpose-Built Autonomous Manufacturing Facility",
                    "Valer, Ostfold \u2014 2,635 m\u00b2 on 30,000 m\u00b2 site (Norbygg)")

    # Left column — bullet points
    items = [
        "5 MAZAK 5-axis CNC machines at production start (August 2027)",
        "Scalable to 20+ CNC machines by 2030",
        "Wenzel CMM measurement room \u2014 precision quality control",
        "Fully integrated autonomous process control",
        "Digital thread from customer order to final delivery",
        "Kaeser compressed air, Absolent extraction systems",
    ]
    add_bullet_list(slide, MARGIN, Inches(1.4), Inches(5.3), Inches(2.8), items, font_size=12)

    # Right stat boxes
    sx = Inches(6.3)
    # Box 1
    add_rect(slide, sx, Inches(1.4), Inches(3.2), Inches(1.2), fill_color=PANEL_GRAY)
    add_textbox(slide, sx, Inches(1.5), Inches(3.2), Inches(0.5),
                "8,760 hours/year", font_size=26, font_color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, sx, Inches(2.0), Inches(3.2), Inches(0.35),
                "Theoretical capacity per machine (24/7)", font_size=10, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # Box 2
    add_rect(slide, sx, Inches(2.8), Inches(3.2), Inches(1.2), fill_color=PANEL_GRAY)
    add_textbox(slide, sx, Inches(2.9), Inches(3.2), Inches(0.5),
                "0.8 FTE / CNC", font_size=26, font_color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, sx, Inches(3.4), Inches(3.2), Inches(0.35),
                "vs industry average 2.5 \u2014 autonomous advantage", font_size=10, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # Footer note
    add_textbox(slide, MARGIN, Inches(4.5), CONTENT_WIDTH, Inches(0.3),
                "Facility construction Q3 2026 \u2013 Q1 2027  |  Machine delivery Q2 2027  |  Production start August 2027",
                font_size=10, font_color=GRAY, italic=True)

    add_footer(slide)


def build_slide_09_certifications(prs):
    """CERTIFICATIONS — Timeline + segment alignment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Quality Certifications: Phased Roadmap",
                    "ISO 9001 \u2192 AS9100 \u2192 AQAP 2110 \u2014 aligned with customer requirements")

    # Timeline — three milestones
    line_y = Inches(2.1)
    add_rect(slide, Inches(1.0), line_y + Inches(0.15), Inches(8.0), Inches(0.03), fill_color=DIVIDER_GRAY)

    milestones = [
        (Inches(1.5), "Q4 2027", "ISO 9001", "Foundation quality\nmanagement system"),
        (Inches(4.5), "Q3-Q4 2028", "AS9100 Rev D", "Aerospace quality\nstandard"),
        (Inches(7.5), "Q3-Q4 2028", "AQAP 2110", "NATO defense quality\nassurance"),
    ]

    for mx, date, cert, desc in milestones:
        # Circle marker
        add_circle(slide, mx - Inches(0.15), line_y, Inches(0.35), fill_color=RED)
        add_textbox(slide, mx - Inches(0.15), line_y + Inches(0.02), Inches(0.35), Inches(0.35),
                    "\u2713", font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Date above
        add_textbox(slide, mx - Inches(0.8), line_y - Inches(0.4), Inches(1.6), Inches(0.3),
                    date, font_size=12, font_color=RED, bold=True, alignment=PP_ALIGN.CENTER)
        # Cert name below
        add_textbox(slide, mx - Inches(0.8), line_y + Inches(0.45), Inches(1.6), Inches(0.3),
                    cert, font_size=13, font_color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
        # Description
        add_textbox(slide, mx - Inches(0.8), line_y + Inches(0.75), Inches(1.6), Inches(0.5),
                    desc, font_size=10, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # Segment alignment boxes
    box_y = Inches(3.6)
    box_h = Inches(0.6)
    boxes = [
        (MARGIN, "Defense: AQAP 2110 qualification target Q3-Q4 2028"),
        (Inches(3.5), "Energy: ISO 9001 + NORSOK alignment from Q4 2027"),
        (Inches(6.7), "Maritime: ISO 9001 sufficient from Q4 2027"),
    ]
    for bx, text in boxes:
        add_rect(slide, bx, box_y, Inches(2.9), box_h, fill_color=PANEL_GRAY)
        add_textbox(slide, bx + Inches(0.1), box_y + Inches(0.1), Inches(2.7), box_h - Inches(0.2),
                    text, font_size=10, font_color=BLACK, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def build_slide_10_capabilities(prs):
    """MACHINE CAPABILITIES — Table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Machine Capabilities & Specifications",
                    "MAZAK 5-axis CNC portfolio \u2014 precision, versatility, reliability")

    rows = [
        ["Capability", "Specification"],
        ["Machine type", "MAZAK 5-axis CNC machining centers"],
        ["Material compatibility", "Aluminum, steel, stainless, titanium, Inconel, composites"],
        ["Precision tolerance", "\u00b10.005 mm typical, \u00b10.002 mm achievable"],
        ["Part size range", "Small precision components to 1,200 mm envelope"],
        ["Surface finish", "Aerospace-grade surface finishes, complex geometries"],
        ["Measurement", "Wenzel CMM in-house \u2014 verifiable quality traceability"],
        ["Capacity (5 CNC)", "43,800 hours/year theoretical (8,760 \u00d7 5)"],
        ["Capacity at 60% util.", "26,280 productive hours/year delivered"],
    ]

    create_table(slide, MARGIN, Inches(1.35), Inches(9.0), rows,
                 col_widths=[Inches(2.5), Inches(6.5)])

    add_textbox(slide, MARGIN, Inches(4.7), CONTENT_WIDTH, Inches(0.3),
                "Custom tooling, fixtures, and process development available. Specifications subject to final equipment selection.",
                font_size=9, font_color=GRAY, italic=True)

    add_footer(slide)


def build_slide_11_pricing(prs):
    """PRICING MODEL — Standard + Customer Program."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Pricing: Transparent & Partnership-Driven",
                    "Aligned incentives for volume commitments")

    # Left column — Standard pricing
    add_textbox(slide, MARGIN, Inches(1.4), Inches(4.0), Inches(0.35),
                "STANDARD PRICING", font_size=16, font_color=BLACK, bold=True)

    # Red badge with price (white text on red)
    create_stat_badge(slide, Inches(0.8), Inches(1.9), Inches(3.0), Inches(0.7),
                      "3,000 NOK / hour", "Conservative floor \u2014 competitive with industry 2,500-3,500 NOK range")

    items_left = [
        "Transparent hourly billing",
        "No hidden fees or setup charges for standard work",
        "Predictable cost planning for customers",
    ]
    add_bullet_list(slide, MARGIN, Inches(3.2), Inches(4.2), Inches(1.5), items_left, font_size=11)

    # Right column — Customer Program (callout style)
    cp_left = Inches(5.3)
    cp_width = Inches(4.2)
    add_rect(slide, cp_left, Inches(1.4), Inches(0.04), Inches(3.2), fill_color=RED)
    add_rect(slide, cp_left + Inches(0.04), Inches(1.4), cp_width - Inches(0.04), Inches(3.2), fill_color=PANEL_GRAY)

    add_textbox(slide, cp_left + Inches(0.2), Inches(1.5), cp_width - Inches(0.4), Inches(0.35),
                "CUSTOMER PROGRAM", font_size=16, font_color=BLACK, bold=True)
    add_textbox(slide, cp_left + Inches(0.2), Inches(1.9), cp_width - Inches(0.4), Inches(0.35),
                "For strategic customers committing >100,000 hours/year:",
                font_size=11, font_color=BLACK, bold=True)

    cp_items = [
        "50/50 profit-sharing above 45% utilization threshold",
        "Incentivizes volume, guarantees capacity, aligns success",
        "Customers get cost reduction at scale",
        "Aurelian secures committed production volume",
    ]
    add_bullet_list(slide, cp_left + Inches(0.2), Inches(2.4), cp_width - Inches(0.4), Inches(1.8),
                    cp_items, font_size=11, spacing_pt=4)

    add_textbox(slide, cp_left + Inches(0.2), Inches(4.0), cp_width - Inches(0.4), Inches(0.4),
                "Qualification: Defense Tier 1, Energy majors, and other strategic accounts",
                font_size=9, font_color=GRAY, italic=True)

    add_footer(slide)


def build_slide_12_reliability(prs):
    """DELIVERY RELIABILITY — Four benefit cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "What You Get: Absolute Delivery Reliability",
                    "Verifiable quality, predictable timelines, digital transparency")

    card_w = Inches(4.2)
    card_h = Inches(1.45)
    gap_x = Inches(0.5)
    gap_y = Inches(0.2)

    cards = [
        ("Verifiable Quality Traceability",
         "Digital thread from quote to delivery \u2014 every part traceable, every measurement recorded and auditable"),
        ("Predictable Capacity",
         "Guaranteed machine hours, autonomous scheduling eliminates manual bottlenecks and shift dependencies"),
        ("European Trusted Framework",
         "Norwegian-based, GDPR-compliant, aligned with EU defense and energy security priorities"),
        ("Scalable Partnership",
         "Growth path from prototype runs to serial production \u2014 we scale capacity with your program needs"),
    ]

    for i, (title, body) in enumerate(cards):
        row = i // 2
        col = i % 2
        x = MARGIN + col * (card_w + gap_x)
        y = Inches(1.45) + row * (card_h + gap_y)
        create_card(slide, x, y, card_w, card_h, title, body, title_size=13, body_size=11)

    add_footer(slide)


def build_slide_13_timeline(prs):
    """PRODUCTION TIMELINE."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Production Timeline: August 2027",
                    "Key milestones from funding to full operations")

    # Horizontal timeline line
    line_y = Inches(2.5)
    add_rect(slide, Inches(0.5), line_y, Inches(9.0), Inches(0.03), fill_color=DIVIDER_GRAY)

    milestones = [
        (Inches(0.7), "Q1-Q2\n2026", "Seed Funding\nClose", False),
        (Inches(2.2), "Q3\n2026", "Facility Build\nStarts", False),
        (Inches(3.7), "Q1\n2027", "Facility\nComplete", False),
        (Inches(5.2), "Q2\n2027", "Machines\nDelivered", False),
        (Inches(6.7), "AUG\n2027", "PRODUCTION\nSTART", True),  # Highlighted
        (Inches(8.2), "Q4\n2027", "ISO 9001\nCertified", False),
    ]

    for mx, date, label, highlight in milestones:
        circle_color = RED
        circle_size = Inches(0.3) if not highlight else Inches(0.4)
        offset = Inches(0.0) if not highlight else Inches(-0.05)

        add_circle(slide, mx, line_y - Inches(0.13) + offset, circle_size, fill_color=circle_color)

        # Date above
        date_color = RED if highlight else GRAY
        date_size = 11 if highlight else 10
        add_textbox(slide, mx - Inches(0.4), line_y - Inches(0.85), Inches(1.1), Inches(0.6),
                    date, font_size=date_size, font_color=date_color, bold=True, alignment=PP_ALIGN.CENTER)

        # Label below
        label_color = RED if highlight else BLACK
        label_size = 11 if highlight else 10
        add_textbox(slide, mx - Inches(0.4), line_y + Inches(0.3), Inches(1.1), Inches(0.6),
                    label, font_size=label_size, font_color=label_color,
                    bold=highlight, alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(slide, MARGIN, Inches(4.2), CONTENT_WIDTH, Inches(0.5),
                "AS9100 / AQAP 2110 qualification: Q3-Q4 2028  |  Serie A expansion to 20 CNC: 2029-2030",
                font_size=10, font_color=GRAY, italic=True, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def build_slide_14_reference(prs):
    """REFERENCE CUSTOMER — Physical Robotics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Reference Customer: Physical Robotics",
                    "Letter of Intent signed \u2014 precision components for autonomous systems")

    # Left content
    items = [
        "Norwegian autonomous systems developer",
        "Requirement: High-precision aluminum and steel components",
        "Volume commitment supporting serial production ramp-up from 2027",
        "Early adopter of Customer Program pricing model",
    ]
    add_bullet_list(slide, MARGIN, Inches(1.4), Inches(5.3), Inches(2.0), items, font_size=12)

    # Callout quote
    add_callout_box(slide, MARGIN, Inches(3.5), Inches(5.3), Inches(0.8),
                    "First reference customer validates both technical fit and commercial model for autonomous manufacturing supply chains.",
                    font_size=11)

    # Right dark panel
    panel_left = Inches(6.3)
    panel_width = Inches(3.2)
    add_rect(slide, panel_left, Inches(1.4), panel_width, Inches(2.8), fill_color=BLACK)

    add_textbox(slide, panel_left + Inches(0.2), Inches(1.55), Inches(2.8), Inches(0.35),
                "WHAT THIS PROVES", font_size=14, font_color=WHITE, bold=True)

    proof_items = [
        "Real customer validation",
        "Technical fit confirmed",
        "Commercial model proven",
        "Reference for defense/energy discussions",
    ]
    for i, item in enumerate(proof_items):
        y = Inches(2.05) + i * Inches(0.45)
        add_textbox(slide, panel_left + Inches(0.2), y, Inches(2.8), Inches(0.35),
                    f"\u2713  {item}", font_size=12, font_color=WHITE)

    add_footer(slide)


def build_slide_15_competitive(prs):
    """COMPETITIVE COMPARISON — Table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "Why Aurelian vs. Traditional Job Shops",
                    "Greenfield advantage, autonomous operations, aligned incentives")

    rows = [
        ["Factor", "Traditional Job Shops", "Aurelian Manufacturing"],
        ["Utilization", "20-45% (observed)", "60-65% (design target)"],
        ["Staffing ratio", "~2.5 FTE/CNC", "0.8 FTE/CNC (sub-linear)"],
        ["Operating hours", "Single shift, 5-day week", "24/7 autonomous operations"],
        ["Quality traceability", "Manual records, variable", "Digital thread, verifiable"],
        ["Capacity predictability", "Labor-dependent, volatile", "Guaranteed, contractual"],
        ["Pricing model", "Transactional, variable", "Transparent + partnership"],
        ["Certifications", "Varies, often incomplete", "ISO \u2192 AS9100 \u2192 AQAP roadmap"],
    ]

    create_table(slide, MARGIN, Inches(1.35), Inches(9.0), rows,
                 col_widths=[Inches(2.2), Inches(3.4), Inches(3.4)])

    add_textbox(slide, MARGIN, Inches(4.55), CONTENT_WIDTH, Inches(0.4),
                "Industry data: observed values (VDR 03.02 CNC Benchmark). Aurelian figures: design targets for greenfield autonomous facility. See \u00a72.2 claims discipline.",
                font_size=9, font_color=GRAY, italic=True)

    add_footer(slide)


def build_slide_16_engagement(prs):
    """ENGAGEMENT PROCESS — Four numbered steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide,
                    "How to Work With Us",
                    "Transparent, phased approach aligned with your procurement cycles")

    steps = [
        ("Initial Consultation", "Technical requirements, capacity planning, certification timeline alignment"),
        ("Technical Validation", "Sample part review, process development, quality requirements mapping"),
        ("Commercial Agreement", "Standard pricing or Customer Program, volume commitment, SLA definition"),
        ("Production Onboarding", "Digital thread integration, first article inspection, serial production ramp-up"),
    ]

    # Left column — steps
    for i, (title, desc) in enumerate(steps):
        y = Inches(1.45) + i * Inches(0.8)
        add_numbered_step(slide, MARGIN, y, i + 1, title, desc)

    # Right — timeline by tier
    add_rect(slide, Inches(5.8), Inches(1.5), Inches(3.7), Inches(2.8), fill_color=PANEL_GRAY)
    add_textbox(slide, Inches(6.0), Inches(1.6), Inches(3.3), Inches(0.3),
                "TYPICAL TIMELINES", font_size=13, font_color=BLACK, bold=True)

    tiers = [
        ("Tier 1: Defense", "12-18 months", RED),
        ("Tier 2: Energy", "6-12 months", BLACK),
        ("Tier 3: Maritime/Infra", "3-6 months", GRAY),
    ]
    for i, (tier, timeline, color) in enumerate(tiers):
        y = Inches(2.1) + i * Inches(0.65)
        add_textbox(slide, Inches(6.0), y, Inches(2.0), Inches(0.25),
                    tier, font_size=11, font_color=color, bold=True)
        add_textbox(slide, Inches(8.0), y, Inches(1.5), Inches(0.25),
                    timeline, font_size=11, font_color=BLACK)

    add_textbox(slide, Inches(6.0), Inches(3.6), Inches(3.3), Inches(0.4),
                "We align our engagement process with your internal procurement and qualification requirements.",
                font_size=9, font_color=GRAY, italic=True)

    add_footer(slide)


def build_slide_17_references(prs):
    """REFERENCES & CONTACT."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "References & Contact")

    # Left — source documents table
    add_textbox(slide, MARGIN, Inches(1.1), Inches(5.0), Inches(0.3),
                "SOURCE DOCUMENTS", font_size=13, font_color=BLACK, bold=True)

    rows = [
        ["VDR ref", "Document"],
        ["02.04", "02 Economic Tables & Projections (master document)"],
        ["03.02", "CNC Benchmark & Competitive Landscape"],
        ["03.05", "Pricing & Revenue Model"],
        ["03.06", "Market Trends & Projections"],
        ["04.05", "Quality Certification Roadmap"],
        ["04.06", "Production Timeline"],
    ]

    create_table(slide, MARGIN, Inches(1.45), Inches(5.5), rows,
                 col_widths=[Inches(1.0), Inches(4.5)])

    # Right — contact info
    cx = Inches(6.3)
    add_textbox(slide, cx, Inches(1.1), Inches(3.2), Inches(0.3),
                "CONTACT US", font_size=16, font_color=RED, bold=True)

    # Andre
    add_textbox(slide, cx, Inches(1.6), Inches(3.2), Inches(0.25),
                "Andre Tandberg", font_size=13, font_color=BLACK, bold=True)
    add_textbox(slide, cx, Inches(1.85), Inches(3.2), Inches(0.2),
                "CEO & Co-Founder", font_size=11, font_color=GRAY)
    add_textbox(slide, cx, Inches(2.05), Inches(3.2), Inches(0.2),
                "andre@aurelian.no", font_size=11, font_color=BLACK)

    # Tore
    add_textbox(slide, cx, Inches(2.5), Inches(3.2), Inches(0.25),
                "Tore Ausland", font_size=13, font_color=BLACK, bold=True)
    add_textbox(slide, cx, Inches(2.75), Inches(3.2), Inches(0.2),
                "VP Business Development & Co-Founder", font_size=11, font_color=GRAY)
    add_textbox(slide, cx, Inches(2.95), Inches(3.2), Inches(0.2),
                "tore@aurelian.no", font_size=11, font_color=BLACK)

    # Company
    add_textbox(slide, cx, Inches(3.5), Inches(3.2), Inches(0.6),
                "Aurelian Manufacturing AS\nOrg.nr 835 679 632\nValer, Ostfold, Norway",
                font_size=10, font_color=GRAY)

    add_footer(slide)


def build_slide_18_closing(prs, logo_img):
    """CLOSING / CTA — Black background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PURE_BLACK)

    # Logo (small)
    logo_img.seek(0)
    slide.shapes.add_picture(logo_img, Inches(4.5), Inches(0.3), Inches(1.0), Inches(1.0))

    # Heading
    add_textbox(slide, Inches(1.0), Inches(1.4), Inches(8.0), Inches(0.7),
                "PARTNER WITH AURELIAN", font_size=36, font_color=WHITE,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(slide, Inches(1.0), Inches(2.1), Inches(8.0), Inches(0.4),
                "Controllable Capacity. Absolute Reliability. Aligned Success.",
                font_size=18, font_color=RED, alignment=PP_ALIGN.CENTER)

    # Three value boxes
    box_w = Inches(2.5)
    box_h = Inches(1.0)
    box_y = Inches(2.8)
    gap = Inches(0.5)
    start_x = Inches(1.0)

    values = [
        ("August 2027", "Production start"),
        ("3,000 NOK / hour", "Standard rate"),
        ("50/50 Partnership", "Customer Program"),
    ]

    for i, (val, label) in enumerate(values):
        x = start_x + i * (box_w + gap)
        add_rect(slide, x, box_y, box_w, box_h, fill_color=BLACK)
        # Border
        shape = add_rect(slide, x, box_y, box_w, box_h)
        shape.line.fill.solid()
        shape.line.color.rgb = GRAY
        shape.line.width = Pt(0.5)

        add_textbox(slide, x, box_y + Inches(0.1), box_w, Inches(0.45),
                    val, font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, box_y + Inches(0.55), box_w, Inches(0.3),
                    label, font_size=10, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # CTA
    add_textbox(slide, Inches(1.0), Inches(4.15), Inches(8.0), Inches(0.4),
                "Contact us today to explore partnership opportunities",
                font_size=14, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(4.5), Inches(8.0), Inches(0.3),
                "andre@aurelian.no  |  tore@aurelian.no",
                font_size=13, font_color=RED, alignment=PP_ALIGN.CENTER)

    add_footer(slide, bg_is_dark=True)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    print("Aurelian Customer Presentation Generator")
    print("=" * 50)

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Load assets
    print("Loading assets...")
    logo = load_b64_image("logo_white_b64.txt")
    icon_gears = load_b64_image("icon_gears_b64.txt")
    icon_robot = load_b64_image("icon_robot_b64.txt")
    icon_check = load_b64_image("icon_check_b64.txt")
    # icon_oil = load_b64_image("icon_oil_b64.txt")  # Available but not used

    # Build all 18 slides
    print("Building slides...")

    print("  Slide  1: Cover")
    build_slide_01_cover(prs, logo)

    print("  Slide  2: The Problem")
    build_slide_02_problem(prs)

    print("  Slide  3: Our Solution")
    build_slide_03_solution(prs, icon_gears, icon_robot, icon_check)

    print("  Slide  4: Customer Segments")
    build_slide_04_segments(prs)

    print("  Slide  5: Defense Drivers")
    build_slide_05_defense(prs)

    print("  Slide  6: Energy Drivers")
    build_slide_06_energy(prs)

    print("  Slide  7: Maritime Drivers")
    build_slide_07_maritime(prs)

    print("  Slide  8: Facility")
    build_slide_08_facility(prs)

    print("  Slide  9: Certifications")
    build_slide_09_certifications(prs)

    print("  Slide 10: Machine Capabilities")
    build_slide_10_capabilities(prs)

    print("  Slide 11: Pricing Model")
    build_slide_11_pricing(prs)

    print("  Slide 12: Delivery Reliability")
    build_slide_12_reliability(prs)

    print("  Slide 13: Production Timeline")
    build_slide_13_timeline(prs)

    print("  Slide 14: Reference Customer")
    build_slide_14_reference(prs)

    print("  Slide 15: Competitive Comparison")
    build_slide_15_competitive(prs)

    print("  Slide 16: Engagement Process")
    build_slide_16_engagement(prs)

    print("  Slide 17: References & Contact")
    build_slide_17_references(prs)

    print("  Slide 18: Closing")
    build_slide_18_closing(prs, logo)

    # Save
    print(f"\nSaving to: {OUTPUT_FILE}")
    prs.save(str(OUTPUT_FILE))
    print("Done! 18 slides generated successfully.")
    print(f"\nFile: {OUTPUT_FILE.name}")
    print(f"Size: {OUTPUT_FILE.stat().st_size / 1024:.0f} KB")


if __name__ == "__main__":
    main()
