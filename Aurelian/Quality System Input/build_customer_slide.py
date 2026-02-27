"""
Aurelian Manufacturing — Customer-Facing Summary Slide
Single world-class slide: the main message to get customers onboard.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (design manual) ──
RED = RGBColor(0xF5, 0x05, 0x37)
BLACK = RGBColor(0x2B, 0x2B, 0x2B)
PURE_BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_MED = RGBColor(0x6B, 0x6B, 0x6B)
GRAY_PANEL = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_DIV = RGBColor(0xD3, 0xD3, 0xD3)
DARK_CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_text(slide, left, top, width, height, text, font_size=14,
             color=WHITE, bold=False, italic=False,
             alignment=PP_ALIGN.LEFT, font_name=FONT, line_spacing=None):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox, tf


def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=FONT, line_spacing=None):
    """Add a text box with multiple paragraphs (list of (text, color, bold, size) tuples)."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, clr, bld, sz = item, color, bold, font_size
        else:
            txt = item[0]
            clr = item[1] if len(item) > 1 else color
            bld = item[2] if len(item) > 2 else bold
            sz = item[3] if len(item) > 3 else font_size

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.alignment = alignment
        if line_spacing:
            p.line_spacing = Pt(line_spacing)

    return txBox, tf


def add_rect(slide, left, top, width, height, fill_color,
             text="", font_size=11, text_color=WHITE, bold=False,
             border_color=None, border_width=None, radius=0.05):
    """Add a rounded rectangle with optional text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_rect_multiline(slide, left, top, width, height, fill_color,
                       lines, border_color=None, border_width=None,
                       radius=0.05):
    """Add a rounded rectangle with multiple text lines."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)

    for i, (txt, clr, bld, sz, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = FONT
        p.alignment = align

    return shape


# ════════════════════════════════════════════════════════════════
# THE SLIDE — BLACK BACKGROUND, CUSTOMER-FACING SUMMARY
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = PURE_BLACK

# ── Top: Company identifier ──
add_text(slide, 0, 0.35, 13.333, 0.35,
         "AURELIAN MANUFACTURING",
         font_size=12, color=GRAY_MED, bold=True,
         alignment=PP_ALIGN.CENTER)

# ── Hero headline ──
add_text(slide, 1.0, 1.0, 11.333, 0.7,
         "Your Product. Our Digital Thread.",
         font_size=44, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

# ── Red accent bar ──
add_rect(slide, 5.667, 1.85, 2.0, 0.04, RED, radius=0.0)

# ── Value statement ──
add_text(slide, 1.5, 2.1, 10.333, 0.55,
         "Every part ships with a complete, digitally assembled Manufacturing Record Book —\n"
         "from raw material certificate to final inspection. Defence and Energy grade.",
         font_size=16, color=GRAY_MED, bold=False, italic=True,
         alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# FOUR VALUE PILLARS — arranged horizontally
# ════════════════════════════════════════════════════════════════
pillar_y = 3.05
pillar_h = 2.8
pillar_w = 2.65
gap = 0.28
start_x = 0.85  # center 4 pillars across 13.333"

pillars = [
    {
        "icon": "⬡",  # hexagon — traceability
        "number": "100%",
        "title": "TRACEABILITY",
        "bullets": [
            "Material → Process → Inspection → Shipment",
            "Real-time document assembly",
            "Customer-specific SDRL compliance",
            "AS9102 FAIR, EN 10204, NORSOK M-630",
        ]
    },
    {
        "icon": "◈",  # diamond — quality
        "number": "5-LAYER",
        "title": "VALIDATION",
        "bullets": [
            "Existence → Format → Completeness",
            "→ Compliance → Traceability",
            "Automated Material Gate (5-point cert check)",
            "Zero manual transcription errors",
        ]
    },
    {
        "icon": "⬢",  # shield — security
        "number": "AES-256",
        "title": "CYBER SECURITY",
        "bullets": [
            "Encryption at rest and in transit",
            "Customer data isolation",
            "Role-based access (6 levels + MFA)",
            "ITAR/EAR physical data separation",
        ]
    },
    {
        "icon": "◉",  # circle — longevity
        "number": "30+ YR",
        "title": "RETENTION",
        "bullets": [
            "Defence-grade document archival",
            "RPO < 24h  •  RTO < 4h",
            "Annual integrity audit (SHA-256)",
            "Technology migration plan built in",
        ]
    },
]

for i, pil in enumerate(pillars):
    px = start_x + i * (pillar_w + gap)

    # Card background — dark charcoal
    add_rect(slide, px, pillar_y, pillar_w, pillar_h,
             DARK_CHARCOAL, border_color=RGBColor(0x35, 0x35, 0x35),
             border_width=1, radius=0.04)

    # Number — white on red badge at top of card
    badge_w = 1.8
    badge_x = px + (pillar_w - badge_w) / 2
    add_rect(slide, badge_x, pillar_y + 0.2, badge_w, 0.5,
             RED, text=pil["number"], font_size=22,
             text_color=WHITE, bold=True, radius=0.06)

    # Title
    add_text(slide, px, pillar_y + 0.85, pillar_w, 0.35,
             pil["title"],
             font_size=13, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

    # Thin red separator
    sep_w = 1.2
    sep_x = px + (pillar_w - sep_w) / 2
    add_rect(slide, sep_x, pillar_y + 1.15, sep_w, 0.025, RED, radius=0.0)

    # Bullet points
    bullet_lines = []
    for j, b in enumerate(pil["bullets"]):
        bullet_lines.append(
            (b, GRAY_MED, False, 10, PP_ALIGN.LEFT)
        )

    add_rect_multiline(
        slide, px + 0.1, pillar_y + 1.25, pillar_w - 0.2, pillar_h - 1.45,
        DARK_CHARCOAL, bullet_lines)

# ════════════════════════════════════════════════════════════════
# BOTTOM — Industry coverage bar
# ════════════════════════════════════════════════════════════════
bar_y = 6.1
add_rect(slide, 0.85, bar_y, 11.633, 0.55, RED, radius=0.03)

# Industry text inside the red bar
txBox = slide.shapes.add_textbox(
    Inches(0.85), Inches(bar_y), Inches(11.633), Inches(0.55))
tf = txBox.text_frame
tf.word_wrap = True
tf.margin_top = Pt(8)
tf.margin_bottom = Pt(8)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
# Build runs for mixed formatting
run1 = p.add_run()
run1.text = "DEFENCE  "
run1.font.size = Pt(13)
run1.font.color.rgb = WHITE
run1.font.bold = True
run1.font.name = FONT

run2 = p.add_run()
run2.text = "AS9100 • AQAP 2110 • NADCAP • AS9102"
run2.font.size = Pt(11)
run2.font.color.rgb = WHITE
run2.font.bold = False
run2.font.name = FONT

run3 = p.add_run()
run3.text = "      │      "
run3.font.size = Pt(13)
run3.font.color.rgb = WHITE
run3.font.bold = False
run3.font.name = FONT

run4 = p.add_run()
run4.text = "ENERGY & MARITIME  "
run4.font.size = Pt(13)
run4.font.color.rgb = WHITE
run4.font.bold = True
run4.font.name = FONT

run5 = p.add_run()
run5.text = "NORSOK M-630 • M-650 • API • DNV"
run5.font.size = Pt(11)
run5.font.color.rgb = WHITE
run5.font.bold = False
run5.font.name = FONT

# ── Footer ──
add_text(slide, 0.5, 6.95, 4, 0.3,
         "CONFIDENTIAL",
         font_size=10, color=GRAY_MED)

add_text(slide, 5, 6.95, 7.8, 0.3,
         "aurelian.no  •  Valer, Norway  •  February 2026",
         font_size=10, color=GRAY_MED,
         alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SPEAKER NOTES — Abbreviation glossary + retention justification
# ════════════════════════════════════════════════════════════════
notes_slide = slide.notes_slide
notes_tf = notes_slide.notes_text_frame
notes_tf.text = ""  # clear default

notes_text = """ABBREVIATIONS USED ON THIS SLIDE

TRACEABILITY PILLAR:
- SDRL = Supplier Document Requirements List — the customer's list of exactly which documents must be delivered with each part (certificates, test reports, inspection records, etc.). Every customer has a different SDRL. Our system parses these automatically.
- AS9102 = Aerospace standard for First Article Inspection Report (FAIR). Defines the format for proving a new part meets all drawing requirements. Three forms: Form 1 (Part Number Accountability), Form 2 (Product Accountability — Raw Material), Form 3 (Characteristic Accountability — Dimensions).
- FAIR = First Article Inspection Report — the formal proof that a manufactured part meets all design requirements. Required before series production begins for defence/aerospace customers.
- EN 10204 = European standard for metallic products inspection documents. Type 3.1 = Inspection certificate issued by manufacturer. Type 3.2 = Inspection certificate validated by both manufacturer and independent inspector. Defence and energy customers typically require 3.1 or 3.2 level material certificates.
- NORSOK M-630 = Material Data Sheet standard for the Norwegian petroleum industry. Defines documentation requirements for materials and equipment used on the Norwegian Continental Shelf. Equivalent in purpose to EN 10204 but specific to Oil & Gas.

VALIDATION PILLAR:
- 5-Layer Validation = Our progressive validation model: Layer 1 (Existence — is the document there?), Layer 2 (Format — is it the right type/format?), Layer 3 (Completeness — are all required fields filled?), Layer 4 (Compliance — does it meet the specific standard requirements?), Layer 5 (Traceability — can we trace it back through the full chain?).
- Material Gate = The automated checkpoint where incoming material certificates are verified against 5 criteria before material is released to production: (1) Certificate type matches requirement, (2) Heat/lot number matches delivery note, (3) Chemical composition within spec, (4) Mechanical properties within spec, (5) Supplier is on approved supplier list.

CYBER SECURITY PILLAR:
- AES-256 = Advanced Encryption Standard with 256-bit key length. Military-grade encryption used to protect stored data. Same standard used by NATO, NSA, and Norwegian defence agencies.
- MFA = Multi-Factor Authentication — requires two or more verification methods to access the system (e.g., password + hardware token or mobile app). Prevents unauthorized access even if a password is compromised.
- RBAC = Role-Based Access Control — access permissions are assigned to roles (e.g., Quality Manager, QA Engineer, Production Operator) rather than individual users. Ensures people can only see and do what their role requires.
- ITAR = International Traffic in Arms Regulations (US). Controls export of defence articles and services. Requires that ITAR-controlled technical data is physically and logically isolated from non-ITAR data, with access restricted to authorized persons.
- EAR = Export Administration Regulations (US). Controls export of dual-use items (civilian/military). Less restrictive than ITAR but still requires data handling controls.

RETENTION PILLAR:
- RPO = Recovery Point Objective — maximum acceptable amount of data loss measured in time. RPO < 24h means we guarantee no more than 24 hours of data can be lost in a disaster scenario (daily backups).
- RTO = Recovery Time Objective — maximum acceptable time to restore the system after a failure. RTO < 4h means full operational recovery within 4 business hours.
- SHA-256 = Secure Hash Algorithm producing a 256-bit digital fingerprint. Used to verify document integrity — if a single bit changes in a document, the hash changes completely. We store SHA-256 checksums for every document and audit 10% of the archive annually.

INDUSTRY STANDARDS BAR:
- AS9100 = Quality Management System standard for aviation, space, and defence. Built on ISO 9001 with additional requirements for product safety, risk management, and configuration control.
- AQAP 2110 = Allied Quality Assurance Publication 2110 (NATO). Quality requirements for defence suppliers. Required by Norwegian Defence Materiel Agency (FMA) and NATO procurement.
- NADCAP = National Aerospace and Defense Contractors Accreditation Program. Third-party accreditation for special processes (welding, heat treatment, NDT, surface treatment). Managed by the Performance Review Institute (PRI).
- NORSOK M-650 = Qualification of manufacturers of special materials. Defines requirements for suppliers manufacturing materials for use on the Norwegian Continental Shelf.
- API = American Petroleum Institute. Publishes standards for oil and gas equipment manufacturing (e.g., API 6A for wellhead equipment, API 5CT for casing).
- DNV = Det Norske Veritas (now DNV GL). Norwegian classification society. Provides third-party inspection, certification, and classification for maritime, oil & gas, and energy industries.

30+ YEAR RETENTION — INDUSTRY JUSTIFICATION:
This is not an ambitious target — it is the industry minimum for our target markets:
- Defence/Aerospace: AS9100 and AQAP 2110 require retention for the lifetime of the product plus disposal period. Fighter aircraft (F-35: 50+ year service life), naval vessels (30-40 years), and missile systems all require documentation to be available for the full operational life. NATO contracts routinely specify 30+ years.
- Oil & Gas: Equipment on the Norwegian Continental Shelf has a design life of 25-30 years. Equinor and other operators require documentation for the full operational life plus decommissioning, which typically extends to 35-45 years total.
- Norwegian Law: Arkivlova (Archives Act) governs public records. Bokforingsloven requires 5-year minimum for financial records. Quality records for critical infrastructure fall under significantly longer requirements.
- Our approach: PDF/A-3 as baseline format (ISO 19005-3, designed for long-term archival), technology migration plan reviewed every 5 years to prevent format obsolescence, no vendor lock-in on archive storage. This is what differentiates us from competitors who simply "keep files on a server."

WHY THIS APPLIES TO AURELIAN AS A PARTS SUPPLIER (NOT ONLY OEMs):
Retention obligations flow down the supply chain. Even though Aurelian machines parts for assembly at the customer's site, the full retention requirement applies:
- Manufacturer of record: Aurelian's name is on the Certificate of Conformance (CoC). We received the raw material, verified the certs, machined the part, inspected it, and shipped it. We own the documentation for that part.
- Contractual flow-down: AS9100D 7.5.3.2 requires quality records to be retained for the period defined by the customer or regulatory requirements. These flow down from the prime contractor (Kongsberg, Nammo, Equinor) to every sub-tier supplier. The customer dictates the retention period in the PO or framework agreement — Aurelian does not get to choose.
- Failure investigation liability: If a part fails in 20 years and the root cause investigation traces back to a material defect or process deviation, the investigation follows the traceability chain back to Aurelian's MRB — the material certificate, dimensional inspection, process parameters, and operator records. These must be retrievable.
- Defence (AQAP 2110 4.2.4): The supplier shall maintain records for the period specified in the contract. NATO contracts routinely specify lifetime of the system (30-50+ years). This applies to every supplier in the chain, not just the prime.
- Oil & Gas (NORSOK Z-001 / DNV): The operator (e.g., Equinor) requires documentation for the full field life. DNV classification rules require that material certificates for pressure-retaining parts are traceable for the full service life. This flows down to the parts manufacturer.
- Competitive advantage: Most small-to-medium CNC shops keep paper binders or basic file servers. When a customer asks "produce the MRB for part X from 2029" in the year 2050, Aurelian can retrieve it instantly with digital search. Traditional shops often cannot — records are lost, water-damaged, in retired file formats, or the company has changed ownership. This is a real differentiator in supplier qualification.
"""

notes_tf.text = notes_text

# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
import os
out_dir = r"C:\Users\mrntau\OneDrive - Mosseregionens Næringsutvikling AS\Skrivebord\Aurelian\Quality System Input"
out_path = os.path.join(out_dir, "Aurelian_Customer_Summary_Slide_V3.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
