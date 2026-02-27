"""
Aurelian Manufacturing — Digital Technology Layer
Single PPTX slide per design manual v2.0
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (design manual) ──
RED = RGBColor(0xF5, 0x05, 0x37)
BLACK = RGBColor(0x2B, 0x2B, 0x2B)
PURE_BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_MED = RGBColor(0x6B, 0x6B, 0x6B)
GRAY_PANEL = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_DIV = RGBColor(0xD3, 0xD3, 0xD3)
DARK_CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)
NAVY_CARD = RGBColor(0x25, 0x3B, 0x52)
STEEL_BLUE = RGBColor(0x8E, 0xAE, 0xC5)
GREEN = RGBColor(0x1C, 0x6E, 0x3D)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ════════════════════════════════════════════
# SLIDE 1 — COVER (black)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = PURE_BLACK

def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox, tf

def add_rect(slide, left, top, width, height, fill_color, text="", font_size=11, text_color=WHITE, bold=False, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
    return shape

# Label
add_text(slide, 0, 0.8, 13.333, 0.4, "AURELIAN MANUFACTURING AS", font_size=13, color=GRAY_MED, bold=False, alignment=PP_ALIGN.CENTER)

# Title
add_text(slide, 0, 1.6, 13.333, 0.8, "Digital Technology Layer", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Red bar
add_rect(slide, 5.167, 2.55, 3.0, 0.05, RED)

# Subtitle
add_text(slide, 0, 2.8, 13.333, 0.5, "Integrated Architecture — From Shop Floor to Customer Portal", font_size=22, color=RED, bold=True, italic=True, alignment=PP_ALIGN.CENTER)

# Description
add_text(slide, 2.5, 3.5, 8.333, 0.9,
    "Full digital tech stack designed from the shop floor by the founding team.\nERP  •  MRB Builder  •  OEE Monitoring  •  CNC Integration  •  Cyber Security  •  Redundancy\nDefence (AS9100/AQAP) and Oil & Gas (NORSOK/API) manufacturing.",
    font_size=14, color=GRAY_MED, alignment=PP_ALIGN.CENTER)

# Footer
add_text(slide, 0.5, 6.8, 3, 0.3, "CONFIDENTIAL", font_size=11, color=GRAY_MED)
add_text(slide, 6, 6.8, 7, 0.3, "February 2026 — Partner Review", font_size=11, color=GRAY_MED, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════
# SLIDE 2 — 4-LAYER ARCHITECTURE
# ════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
bg2 = slide2.background
bg2.fill.solid()
bg2.fill.fore_color.rgb = WHITE

# Title
add_text(slide2, 0.5, 0.3, 9, 0.5, "Digital Technology Layer", font_size=36, color=BLACK, bold=True)
add_text(slide2, 0.5, 0.85, 12, 0.4, "4-Layer Architecture — Designed from the Shop Floor, Built In-House", font_size=18, color=RED, bold=True, italic=True)

# ── LAYER 1: External (customers/suppliers) ──
y = 1.5
add_rect(slide2, 0.5, y, 12.333, 1.0, GRAY_PANEL, border_color=GRAY_DIV, border_width=1)
add_text(slide2, 0.7, y + 0.05, 2.5, 0.3, "LAYER 1 — EXTERNAL", font_size=10, color=RED, bold=True)

# Customer boxes
add_rect(slide2, 0.7, y + 0.35, 2.8, 0.55, WHITE, "Customers\nPOs • SDRLs • Pre-shipment • MRB Download", font_size=9, text_color=BLACK, border_color=GRAY_DIV, border_width=1)
add_rect(slide2, 3.7, y + 0.35, 2.8, 0.55, WHITE, "Suppliers\nMaterial Certs • Special Process Certs • Declarations", font_size=9, text_color=BLACK, border_color=GRAY_DIV, border_width=1)
add_rect(slide2, 6.7, y + 0.35, 2.8, 0.55, WHITE, "Third-Party Inspectors\nDNV • Lloyd's • GQAR • Classification", font_size=9, text_color=BLACK, border_color=GRAY_DIV, border_width=1)
add_rect(slide2, 9.7, y + 0.35, 2.9, 0.55, WHITE, "Regulatory & Standards\nAS9100 • AQAP • NORSOK • ITAR/EAR", font_size=9, text_color=BLACK, border_color=GRAY_DIV, border_width=1)

# Arrow
add_text(slide2, 5.5, y + 0.95, 2, 0.3, "▼   REST API / Webhooks   ▼", font_size=9, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# ── LAYER 2: Business Operations (ERP) ──
y2 = 2.85
add_rect(slide2, 0.5, y2, 5.9, 1.15, DARK_CHARCOAL)
add_text(slide2, 0.7, y2 + 0.05, 3, 0.25, "LAYER 2 — ERP (BOUGHT SYSTEM)", font_size=10, color=RED, bold=True)
add_text(slide2, 0.7, y2 + 0.32, 5.5, 0.75,
    "Financial Management  •  Purchasing / SCM  •  Production Orders\n"
    "Inventory / Lot Tracking  •  CRM / Sales  •  HR / Payroll\n"
    "Calibration Records  •  Supplier Approvals  •  Compliance  •  NCR",
    font_size=10, color=STEEL_BLUE)

# Interface connector
add_rect(slide2, 6.6, y2 + 0.15, 0.85, 0.85, RED, "20\nInterfaces\nAM-IF-001", font_size=9, text_color=WHITE, bold=True)
# Arrows
add_text(slide2, 6.25, y2 + 0.35, 0.35, 0.3, "◀▶", font_size=14, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# ── LAYER 2b: MRB Builder (Aurelian-built) ──
add_rect(slide2, 7.65, y2, 5.2, 1.15, RED)
add_text(slide2, 7.85, y2 + 0.05, 4, 0.25, "LAYER 2 — MRB BUILDER + OEE (AURELIAN-BUILT)", font_size=10, color=WHITE, bold=True)
add_text(slide2, 7.85, y2 + 0.32, 4.8, 0.75,
    "Material Gate (PR-008)  •  SDRL Parser (PR-009)\n"
    "Document Validator (PR-010)  •  CoC Generator (PR-011)\n"
    "MRB Assembler (PR-012)  •  Archive Engine (PR-013)\n"
    "OEE Engine (Availability × Performance × Quality)\n"
    "Customer Portal  •  Real-time Dashboards",
    font_size=10, color=WHITE)

# Arrow down
add_text(slide2, 5.5, y2 + 1.1, 2.5, 0.3, "▼   MQTT / OPC-UA / MTConnect   ▼", font_size=9, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# ── LAYER 3: Shop Floor ──
y3 = 4.35
add_rect(slide2, 0.5, y3, 12.333, 0.95, NAVY_CARD)
add_text(slide2, 0.7, y3 + 0.05, 3, 0.25, "LAYER 3 — SHOP FLOOR (DIRECT INTEGRATION)", font_size=10, color=RED, bold=True)

add_rect(slide2, 0.7, y3 + 0.35, 2.2, 0.5, DARK_CHARCOAL, "MAZAK CNC Machines\n5-axis • Multi-tasking", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 3.05, y3 + 0.35, 2.2, 0.5, DARK_CHARCOAL, "CMM / Inspection\nDimensional • Surface", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 5.4, y3 + 0.35, 2.2, 0.5, DARK_CHARCOAL, "Metrology Systems\nTemp • Vibration • Env", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 7.75, y3 + 0.35, 1.8, 0.5, DARK_CHARCOAL, "OEE Data Sources\nUptime • Cycle • Scrap", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 9.7, y3 + 0.35, 1.7, 0.5, DARK_CHARCOAL, "Automated Logistics\nAGV • Tool Mgmt", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 11.55, y3 + 0.35, 1.15, 0.5, DARK_CHARCOAL, "Facility\nAccess • HVAC", font_size=9, text_color=STEEL_BLUE)

# ── LAYER 4: Cyber Security & Redundancy ──
y4 = 5.6
add_rect(slide2, 0.5, y4, 12.333, 1.15, PURE_BLACK)
add_text(slide2, 0.7, y4 + 0.05, 5, 0.25, "LAYER 4 — CYBER SECURITY & REDUNDANCY", font_size=10, color=RED, bold=True)

# Security boxes
sec_y = y4 + 0.35
add_rect(slide2, 0.7, sec_y, 2.0, 0.65, DARK_CHARCOAL, "Encryption\nAES-256 at rest\nTLS 1.3 in transit", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 2.85, sec_y, 2.0, 0.65, DARK_CHARCOAL, "Access Control\nRBAC • MFA\nAudit trail on all actions", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 5.0, sec_y, 2.0, 0.65, DARK_CHARCOAL, "Data Isolation\nCustomer data segregation\nEEA data sovereignty", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 7.15, sec_y, 2.0, 0.65, DARK_CHARCOAL, "Redundancy\nReal-time replication\nOff-site encrypted backup", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 9.3, sec_y, 2.0, 0.65, DARK_CHARCOAL, "Integrity\nSHA-256 checksums\nAnnual 10% sample audit", font_size=9, text_color=STEEL_BLUE)
add_rect(slide2, 11.45, sec_y, 1.2, 0.65, DARK_CHARCOAL, "Recovery\nRPO < 24h\nRTO < 4h", font_size=9, text_color=STEEL_BLUE)

# Footer
add_text(slide2, 0.5, 7.0, 6, 0.3, "CONFIDENTIAL  •  Aurelian Manufacturing AS", font_size=10, color=GRAY_MED)
add_text(slide2, 6, 7.0, 7, 0.3, "Digital Technology Layer — February 2026", font_size=10, color=GRAY_MED, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════
# SLIDE 3 — CYBER SECURITY & REDUNDANCY DETAIL
# ════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
bg3 = slide3.background
bg3.fill.solid()
bg3.fill.fore_color.rgb = WHITE

add_text(slide3, 0.5, 0.3, 9, 0.5, "Cyber Security & Redundancy", font_size=36, color=BLACK, bold=True)
add_text(slide3, 0.5, 0.85, 9, 0.4, "Defence-Grade Protection for 15–30+ Year Document Retention", font_size=18, color=RED, bold=True, italic=True)

# Two columns
col_w = 5.9
gap = 0.533

# LEFT COLUMN — Cyber Security
lx = 0.5
ly = 1.5
add_rect(slide3, lx, ly, col_w, 0.4, BLACK, "CYBER SECURITY", font_size=11, text_color=WHITE, bold=True)

sec_items = [
    ("Encryption at Rest", "AES-256 minimum on all stored data — MRBs, certificates, CoCs, archive. Keys managed via HSM or cloud KMS with automatic rotation."),
    ("Encryption in Transit", "TLS 1.3 for all API communication. Mutual TLS (mTLS) for machine-to-server connections (MQTT/OPC-UA). No unencrypted data leaves the facility network."),
    ("Role-Based Access Control", "6-level access model: System Admin → Quality Manager → Senior QA → QA Engineer → Production → Read-Only. Every action logged with user, timestamp, and IP."),
    ("Multi-Factor Authentication", "MFA required for all system access. Hardware tokens for Quality Manager and above. Session timeout after 15 minutes of inactivity."),
    ("Customer Data Isolation", "Each customer's data logically separated. Defence and Oil & Gas orders cannot cross-reference. ITAR-controlled data physically isolated."),
    ("Network Segmentation", "Shop floor network (CNC/CMM) isolated from business network. MRB Builder in DMZ. Firewall rules enforced per interface (AM-IF-001)."),
]

for i, (title, desc) in enumerate(sec_items):
    item_y = ly + 0.5 + i * 0.8
    add_rect(slide3, lx, item_y, col_w, 0.72, GRAY_PANEL, border_color=GRAY_DIV, border_width=0.5)
    txBox, tf = add_text(slide3, lx + 0.15, item_y + 0.04, col_w - 0.3, 0.2, title, font_size=11, color=RED, bold=True)
    txBox2, tf2 = add_text(slide3, lx + 0.15, item_y + 0.26, col_w - 0.3, 0.42, desc, font_size=9, color=GRAY_MED)

# RIGHT COLUMN — Redundancy
rx = lx + col_w + gap
ry = ly
add_rect(slide3, rx, ry, col_w, 0.4, BLACK, "REDUNDANCY & DISASTER RECOVERY", font_size=11, text_color=WHITE, bold=True)

red_items = [
    ("Real-Time Replication", "Continuous replication to secondary on-site server. Zero data loss for primary server failure. Automatic failover within 60 seconds."),
    ("Daily Incremental Backup", "Every 24 hours to off-site encrypted cloud storage (EEA-based provider, ISO 27001 certified). Integrity verified via SHA-256 checksum comparison."),
    ("Weekly Full Backup", "Complete system snapshot to geographically separate cloud region. Tested quarterly — 5 random MRBs restored and verified end-to-end."),
    ("Annual Archive Snapshot", "Full archive frozen to separate cloud region + optional physical media. Insurance against provider failure. 10% sample integrity audit."),
    ("Recovery Targets", "RPO (Recovery Point Objective): < 24 hours maximum data loss. RTO (Recovery Time Objective): < 4 business hours to full operational recovery."),
    ("Technology Migration Plan", "15–30+ year retention requires format migration. PDF/A-3 as baseline. Migration plan reviewed every 5 years. No vendor lock-in on archive storage."),
]

for i, (title, desc) in enumerate(red_items):
    item_y = ry + 0.5 + i * 0.8
    add_rect(slide3, rx, item_y, col_w, 0.72, GRAY_PANEL, border_color=GRAY_DIV, border_width=0.5)
    txBox, tf = add_text(slide3, rx + 0.15, item_y + 0.04, col_w - 0.3, 0.2, title, font_size=11, color=RED, bold=True)
    txBox2, tf2 = add_text(slide3, rx + 0.15, item_y + 0.26, col_w - 0.3, 0.42, desc, font_size=9, color=GRAY_MED)

# Footer
add_text(slide3, 0.5, 7.0, 6, 0.3, "CONFIDENTIAL  •  Aurelian Manufacturing AS", font_size=10, color=GRAY_MED)
add_text(slide3, 6, 7.0, 7, 0.3, "Digital Technology Layer — February 2026", font_size=10, color=GRAY_MED, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════
# SLIDE 4 — COMPLIANCE & STANDARDS MAPPING
# ════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
bg4 = slide4.background
bg4.fill.solid()
bg4.fill.fore_color.rgb = WHITE

add_text(slide4, 0.5, 0.3, 9, 0.5, "Standards & Compliance Coverage", font_size=36, color=BLACK, bold=True)
add_text(slide4, 0.5, 0.85, 9, 0.4, "Every Layer Mapped to Industry Requirements", font_size=18, color=RED, bold=True, italic=True)

# Callout box
add_rect(slide4, 0.5, 1.45, 12.333, 0.65, GRAY_PANEL, border_color=RED, border_width=2)
add_text(slide4, 0.7, 1.5, 11.9, 0.55,
    "The digital technology layer is not a generic IT system. Every component — from the Material Gate certificate check to the archive integrity audit — "
    "is designed around specific industry standards. The system enforces compliance in real time, not retroactively.",
    font_size=12, color=BLACK)

# Standards mapping table — built with shapes
table_y = 2.35
headers = ["Layer", "Component", "Defence Standards", "Oil & Gas Standards", "Security / Regulatory"]
col_widths = [1.5, 2.5, 2.7, 2.7, 2.933]
col_starts = [0.5]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

# Header row
for i, (header, cw) in enumerate(zip(headers, col_widths)):
    add_rect(slide4, col_starts[i], table_y, cw, 0.35, BLACK, header, font_size=10, text_color=WHITE, bold=True)

rows = [
    ["External", "Customer Interfaces", "AS9100D • AQAP 2110\nAS9102 FAIR", "NORSOK M-650 • M-630\nAPI • ITP (H/W/R)", "ITAR/EAR export control\nGDPR (personal data)"],
    ["ERP", "Business Operations", "AQAP quality records\nGQAR traceability", "NORSOK Z-001 (DFO)\nDNV classification", "Bokføringsloven\nArkivlova"],
    ["MRB Builder", "Quality Documentation\n+ OEE Engine", "AS9102 (FAIR 1-2-3)\nEN 10204 (3.1/3.2)", "NORSOK M-630 (MDS)\nASME Section V (NDT)", "ISO 19005-3 (PDF/A)\nSHA-256 integrity"],
    ["Shop Floor", "Machine Integration\n+ OEE Data", "AS9100D process control\nNADCAP special process", "NORSOK M-650 mfg qual\nASTM test methods", "OPC-UA security model\nNetwork segmentation"],
    ["Security", "Cyber & Redundancy", "AQAP 2110 §4.2\nDefence data handling", "NORSOK Z-001 retention\nOperator requirements", "AES-256 • TLS 1.3\nRBAC • MFA • ISO 27001"],
]

for ri, row in enumerate(rows):
    ry_row = table_y + 0.35 + ri * 0.82
    bg_color = WHITE if ri % 2 == 0 else GRAY_PANEL
    for ci, (cell, cw) in enumerate(zip(row, col_widths)):
        add_rect(slide4, col_starts[ci], ry_row, cw, 0.78, bg_color, border_color=GRAY_DIV, border_width=0.5)
        fc = RED if ci == 0 else (BLACK if ci == 1 else GRAY_MED)
        fb = True if ci <= 1 else False
        fs = 10 if ci <= 1 else 9
        add_text(slide4, col_starts[ci] + 0.1, ry_row + 0.08, cw - 0.2, 0.65, cell, font_size=fs, color=fc, bold=fb)

# Footer
add_text(slide4, 0.5, 7.0, 6, 0.3, "CONFIDENTIAL  •  Aurelian Manufacturing AS", font_size=10, color=GRAY_MED)
add_text(slide4, 6, 7.0, 7, 0.3, "Digital Technology Layer — February 2026", font_size=10, color=GRAY_MED, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════
import os
out_dir = r"C:\Users\mrntau\OneDrive - Mosseregionens Næringsutvikling AS\Skrivebord\Aurelian\Quality System Input"
out_path = os.path.join(out_dir, "Aurelian_Digital_Technology_Layer.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
